"""PPTX generation.

The lightest of the three generators, and the same contract: validated JSON
in, a real deck out. Bullets are capped per slide rather than truncated
silently -- a slide with twenty bullets is unreadable, and quietly dropping
the last fifteen would hide content the agent thought it had delivered, so
the overflow moves to a continuation slide instead.
"""

from __future__ import annotations

from typing import ClassVar

from app.artifacts.content import DeckContent, Slide
from app.artifacts.store import artifact_store
from app.tools.base import ToolContext, ToolResult

MAX_BULLETS_PER_SLIDE = 6


class PptxGenerateTool:
    name = "pptx.generate"
    description = (
        "Build a PowerPoint deck from structured content: a title and a list "
        "of slides, each with a heading and bullets."
    )
    risk_level = "low"
    requires_approval = False
    input_schema: ClassVar[dict] = {
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "subtitle": {"type": "string"},
            "slides": {"type": "array"},
            "filename": {"type": "string"},
        },
        "required": ["title", "slides"],
    }

    def execute(self, args: dict, context: ToolContext) -> ToolResult:
        try:
            content = DeckContent.model_validate(
                {
                    "title": args["title"],
                    "subtitle": args.get("subtitle", ""),
                    "slides": args["slides"],
                }
            )
        except Exception as exc:
            return ToolResult.failed(f"Content did not validate: {exc}")

        filename = args.get("filename") or "deck.pptx"
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"

        record = artifact_store.save(
            task_id=context.task_id,
            artifact_type="pptx",
            filename=filename,
            payload=build_deck(content),
        )
        return ToolResult(
            ok=True,
            data={
                "artifact_id": str(record.id),
                "filename": filename,
                "size_bytes": record.size_bytes,
            },
            detail=f"generated {filename}",
        )


def build_deck(content: DeckContent) -> bytes:
    import io

    from pptx import Presentation

    presentation = Presentation()

    title_layout = presentation.slide_layouts[0]
    opening = presentation.slides.add_slide(title_layout)
    opening.shapes.title.text = content.title
    if len(opening.placeholders) > 1:
        opening.placeholders[1].text = content.subtitle

    bullet_layout = presentation.slide_layouts[1]
    for slide_content in content.slides:
        for part in _paginate(slide_content):
            slide = presentation.slides.add_slide(bullet_layout)
            slide.shapes.title.text = part.heading
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, bullet in enumerate(part.bullets):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = bullet
            if part.notes:
                slide.notes_slide.notes_text_frame.text = part.notes

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _paginate(slide: Slide) -> list[Slide]:
    """Split an overfull slide rather than dropping its tail."""
    if len(slide.bullets) <= MAX_BULLETS_PER_SLIDE:
        return [slide]

    parts: list[Slide] = []
    for start in range(0, len(slide.bullets), MAX_BULLETS_PER_SLIDE):
        chunk = slide.bullets[start : start + MAX_BULLETS_PER_SLIDE]
        heading = (
            slide.heading
            if start == 0
            else f"{slide.heading} (cont. {start // MAX_BULLETS_PER_SLIDE + 1})"
        )
        parts.append(
            Slide(
                heading=heading,
                bullets=chunk,
                notes=slide.notes if start == 0 else "",
            )
        )
    return parts
