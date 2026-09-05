"""Part 04: the tool gateway, the workspace boundary and the generators.

The gateway is the security-relevant piece here. Everything a model can cause
to happen goes through one method, so these tests are mostly about what that
method refuses.
"""

from __future__ import annotations

from typing import ClassVar
from uuid import uuid4

import pytest

from app.artifacts.content import (
    ApprovalNoteContent,
    Citation,
    DeckContent,
    Finding,
    WorkbookContent,
)
from app.artifacts.validator import validate_docx, validate_opens
from app.schemas.shared import Evidence
from app.tools import register_default_tools
from app.tools.base import ToolContext, ToolResult
from app.tools.gateway import _validate, gateway
from app.tools.generators.docx import build_approval_note
from app.tools.generators.pptx import build_deck
from app.tools.generators.xlsx import build_workbook


@pytest.fixture(autouse=True)
def _tools():
    register_default_tools()


@pytest.fixture
def context():
    return ToolContext(
        task_id=uuid4(),
        user_id=uuid4(),
        roles=["ENGINEER"],
        classification="INTERNAL",
        input_file_ids=[],
    )


# --- the gateway ----------------------------------------------------------


def test_every_mvp_tool_is_registered():
    assert set(gateway.names()) >= {
        "knowledge.search",
        "file.read",
        "file.write",
        "ocr.extract",
        "python.execute",
        "docx.generate",
        "xlsx.generate",
        "pptx.generate",
    }


def test_an_unknown_tool_is_refused_not_ignored(context):
    result = gateway.call("shell.exec", {"cmd": "rm -rf /"}, context)
    assert not result.ok
    assert "shell.exec" in result.error


def test_the_catalogue_declares_risk_for_every_tool():
    for entry in gateway.catalogue():
        assert entry["risk_level"] in {"low", "medium", "high"}
        assert entry["description"]
    risks = {entry["name"]: entry["risk_level"] for entry in gateway.catalogue()}
    # Code execution is the one capability that must never be quietly low risk.
    assert risks["python.execute"] == "high"


def test_policy_denial_stops_the_call(context, monkeypatch):
    """A denied tool must not execute, and must say why."""
    from app.integrations import registry

    class DenyAll:
        def check_tool_allowed(self, tool, roles, classification):
            return False, "denied for this test"

        def check_permission(self, **kwargs):
            return True, "ok"

    monkeypatch.setattr(registry, "get_policy", lambda: DenyAll())
    result = gateway.call("file.list", {}, context)
    assert not result.ok
    assert "denied for this test" in result.error


def test_a_tool_that_raises_becomes_a_failed_result(context, monkeypatch):
    """A crashing tool must not take the orchestrator down with it."""

    class Exploding:
        name = "test.explode"
        description = "raises"
        risk_level = "low"
        requires_approval = False
        input_schema: ClassVar[dict] = {
            "type": "object",
            "properties": {},
            "required": [],
        }

        def execute(self, args, ctx):
            raise RuntimeError("boom")

    gateway.register(Exploding())
    result = gateway.call("test.explode", {}, context)
    assert not result.ok
    assert "boom" in result.error


def test_a_tool_call_emits_a_trace_event(context):
    from app.core.events import event_bus

    queue, _ = event_bus.subscribe(context.task_id)
    gateway.call("file.list", {}, context)

    events = []
    while not queue.empty():
        item = queue.get_nowait()
        if item is not None:
            events.append(item.event)
    assert "tool_called" in events
    assert "tool_completed" in events


# --- argument validation --------------------------------------------------


@pytest.mark.parametrize(
    ("args", "schema", "expected"),
    [
        ({}, {"required": ["a"], "properties": {"a": {}}}, "Missing required"),
        ({"b": 1}, {"required": [], "properties": {"a": {}}}, "Unknown argument"),
        (
            {"a": "x"},
            {"required": [], "properties": {"a": {"type": "integer"}}},
            "must be a integer",
        ),
        # bool is a subclass of int; an integer argument must not accept True.
        (
            {"a": True},
            {"required": [], "properties": {"a": {"type": "integer"}}},
            "must be a integer",
        ),
    ],
)
def test_bad_arguments_are_rejected(args, schema, expected):
    assert expected in _validate(args, schema)


def test_valid_arguments_pass():
    schema = {"required": ["a"], "properties": {"a": {"type": "string"}}}
    assert _validate({"a": "ok"}, schema) == ""


def test_a_model_inventing_an_argument_is_refused(context):
    result = gateway.call(
        "knowledge.search", {"query": "V-103", "temperature": 0.9}, context
    )
    assert not result.ok
    assert "temperature" in result.error


# --- the workspace boundary -----------------------------------------------


def test_workspace_refuses_a_path_that_escapes(context):
    from app.tools import workspace

    for hostile in ("../../etc/passwd", "..\\..\\secrets.env", "/etc/passwd"):
        with pytest.raises(workspace.WorkspaceError):
            workspace.resolve(context.task_id, hostile)


def test_workspace_round_trips_a_normal_file(context):
    from app.tools import workspace

    workspace.write(context.task_id, "notes.txt", b"V-103 isolated")
    assert workspace.read(context.task_id, "notes.txt") == b"V-103 isolated"
    assert any(
        entry["name"] == "notes.txt" for entry in workspace.listing(context.task_id)
    )


def test_workspace_refuses_an_oversized_file(context):
    from app.tools import workspace

    with pytest.raises(workspace.WorkspaceError):
        workspace.write(
            context.task_id, "big.bin", b"x" * (workspace.MAX_FILE_BYTES + 1)
        )


def test_file_write_then_read_through_the_gateway(context):
    written = gateway.call(
        "file.write", {"name": "step1.txt", "content": "42"}, context
    )
    assert written.ok

    read = gateway.call("file.read", {"name": "step1.txt"}, context)
    assert read.ok
    assert read.data["text"] == "42"


def test_file_read_will_not_traverse(context):
    result = gateway.call("file.read", {"name": "../../../etc/passwd"}, context)
    assert not result.ok
    assert "outside the task workspace" in result.error


def test_file_read_refuses_a_file_the_task_does_not_own(context):
    """A task may only read the inputs it was created with."""
    result = gateway.call("file.read", {"file_id": str(uuid4())}, context)
    assert not result.ok
    assert "not one of this task's inputs" in result.error


def test_file_read_wants_exactly_one_of_name_or_file_id(context):
    assert not gateway.call("file.read", {}, context).ok
    both = gateway.call(
        "file.read", {"name": "a.txt", "file_id": str(uuid4())}, context
    )
    assert not both.ok


# --- generators -----------------------------------------------------------


def _note() -> ApprovalNoteContent:
    return ApprovalNoteContent(
        title="Approval Note: Pump P-103 Seal Replacement",
        summary="The inspection findings were checked against SOP-204.",
        findings=[
            Finding(
                statement="Isolation was performed at V-103 as required.",
                severity="informational",
                citations=[
                    Citation(document_name="Maintenance SOP.pdf", page=7,
                             section="4.2 Isolation")
                ],
            ),
            Finding(
                statement="The permit was not signed before welding.",
                severity="critical",
                citations=[
                    Citation(document_name="Maintenance SOP.pdf", page=9)
                ],
            ),
        ],
        recommendations=["Re-issue the hot work permit before restart."],
    )


def _evidence() -> list[Evidence]:
    return [
        Evidence(
            document_id=uuid4(),
            document_name="Maintenance SOP.pdf",
            page=7,
            section="4.2 Isolation",
            text="Close valve V-103 and lock it out.",
            score=0.9,
        ),
        Evidence(
            document_id=uuid4(),
            document_name="Maintenance SOP.pdf",
            page=9,
            section="5.1 Permits",
            text="A hot work permit is mandatory before welding.",
            score=0.8,
        ),
    ]


def test_the_docx_generator_produces_an_openable_document():
    payload = build_approval_note(_note())
    report = validate_docx(payload, _note(), _evidence())
    assert report.passed, report.failures


def test_the_note_contains_its_citations():
    import io

    from docx import Document

    document = Document(io.BytesIO(build_approval_note(_note())))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Maintenance SOP.pdf" in text
    assert "p.7" in text
    # Findings are ordered by severity, so the critical one leads.
    assert text.index("CRITICAL") < text.index("INFORMATIONAL")


def test_an_uncited_finding_is_labelled_rather_than_left_looking_supported():
    import io

    from docx import Document

    note = ApprovalNoteContent(
        title="Note",
        summary="s",
        findings=[Finding(statement="A claim with no source.", severity="minor")],
    )
    document = Document(io.BytesIO(build_approval_note(note)))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "not supported by a retrieved document" in text


def test_the_xlsx_generator_produces_an_openable_workbook():
    content = WorkbookContent(
        title="Readings",
        sheets=[
            {
                "name": "Vessels",
                "columns": ["Tag", "Pressure"],
                "rows": [["V-103", 4.2], ["V-104", 3.9]],
            }
        ],
    )
    assert validate_opens(build_workbook(content), "xlsx").passed


def test_a_formula_never_becomes_live_arithmetic():
    """A model emitting "=SUM(...)" must not make the spreadsheet compute."""
    import io

    from openpyxl import load_workbook

    content = WorkbookContent(
        title="Totals",
        sheets=[
            {"name": "S", "columns": ["Value"], "rows": [["=SUM(A1:A9)"], ["12"]]}
        ],
    )
    workbook = load_workbook(io.BytesIO(build_workbook(content)))
    cell = workbook["S"]["A2"].value
    assert not str(cell).startswith("=")


def test_the_pptx_generator_produces_an_openable_deck():
    content = DeckContent(
        title="P&ID Review",
        subtitle="Crude Distillation Unit",
        slides=[{"heading": "Findings", "bullets": ["V-103 isolated"]}],
    )
    assert validate_opens(build_deck(content), "pptx").passed


def test_an_overfull_slide_is_split_not_truncated():
    import io

    from pptx import Presentation

    bullets = [f"point {index}" for index in range(14)]
    content = DeckContent(
        title="T", slides=[{"heading": "Many", "bullets": bullets}]
    )
    deck = Presentation(io.BytesIO(build_deck(content)))

    rendered = " ".join(
        shape.text_frame.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    # Nothing silently disappears.
    for bullet in bullets:
        assert bullet in rendered


def test_a_truncated_file_fails_the_opens_check():
    truncated = build_workbook(
        WorkbookContent(
            title="T", sheets=[{"name": "S", "columns": ["A"], "rows": [["1"]]}]
        )
    )[:200]
    assert not validate_opens(truncated, "xlsx").passed


# --- the validator --------------------------------------------------------


def test_a_citation_to_a_document_that_was_never_retrieved_fails():
    """The most damaging output this system could produce."""
    note = _note()
    note.findings[0].citations = [
        Citation(document_name="Imaginary Standard.pdf", page=3)
    ]
    report = validate_docx(build_approval_note(note), note, _evidence())
    assert not report.passed
    assert any("Imaginary Standard" in failure for failure in report.failures)


def test_a_citation_to_a_page_that_was_never_retrieved_fails():
    note = _note()
    note.findings[0].citations = [
        Citation(document_name="Maintenance SOP.pdf", page=99)
    ]
    report = validate_docx(build_approval_note(note), note, _evidence())
    assert not report.passed
    assert any("p.99" in failure for failure in report.failures)


def test_a_note_made_entirely_of_uncited_claims_fails():
    note = ApprovalNoteContent(
        title="Note",
        summary="s",
        findings=[
            Finding(statement="Claim one.", severity="minor"),
            Finding(statement="Claim two.", severity="minor"),
        ],
    )
    report = validate_docx(build_approval_note(note), note, _evidence())
    assert not report.passed
    assert any("cite nothing" in failure for failure in report.failures)


def test_a_note_with_nothing_in_it_fails():
    note = ApprovalNoteContent(title="Empty", summary="Nothing to report.")
    report = validate_docx(build_approval_note(note), note, [])
    assert not report.passed


def test_a_file_that_will_not_open_fails_everything():
    note = _note()
    report = validate_docx(b"not a docx", note, _evidence())
    assert not report.passed
    assert report.failures


# --- the sandbox ----------------------------------------------------------


def test_a_sandbox_that_cannot_run_says_so_rather_than_reporting_failure(
    context, monkeypatch
):
    """"The code failed" and "nothing ran" mean different things to an agent."""
    from app.sandbox.base import SandboxResult
    from app.tools import python as python_tool

    monkeypatch.setattr(
        python_tool.docker_sandbox,
        "run",
        lambda request: SandboxResult(
            exit_code=-1,
            stdout="",
            stderr="",
            duration_ms=0,
            status="unavailable",
            detail="Docker is not reachable",
        ),
    )
    result = gateway.call("python.execute", {"code": "print(1)"}, context)
    assert not result.ok
    assert "sandbox is unavailable" in result.error
    assert "nothing was run" in result.error


def test_python_execute_rejects_empty_code(context):
    assert not gateway.call("python.execute", {"code": "   "}, context).ok


def test_a_sandbox_result_distinguishes_ran_and_failed_from_never_ran():
    from app.sandbox.base import SandboxResult

    ran_and_failed = SandboxResult(
        exit_code=1, stdout="", stderr="Traceback", duration_ms=5
    )
    assert not ran_and_failed.succeeded
    assert ran_and_failed.status == "ok"

    never_ran = SandboxResult(
        exit_code=-1, stdout="", stderr="", duration_ms=0, status="unavailable"
    )
    assert not never_ran.succeeded


def test_a_hostile_input_filename_cannot_escape_the_sandbox_workspace():
    from app.sandbox.docker_runner import _safe_name

    assert _safe_name("../../etc/passwd") == "passwd"
    assert _safe_name("..\\..\\windows\\system32\\cmd.exe") == "cmd.exe"
    assert _safe_name("") == "input.dat"
    assert _safe_name("report.csv") == "report.csv"


def test_tool_result_failure_helper():
    result = ToolResult.failed("nope", hint="try again")
    assert not result.ok
    assert result.error == "nope"
    assert result.data["hint"] == "try again"
